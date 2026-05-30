"""
RAG Core — no UI dependencies
All document loading, embedding, retrieval, and Claude API logic.
"""

import contextlib
import hashlib
import json
import logging
import os
import time
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, Generator, List, Optional
from dataclasses import dataclass, field

import anthropic
import psycopg2
import psycopg2.pool as pg_pool
import yaml

from langchain_postgres import PGVector
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings
from fastembed import TextEmbedding

import pypdf
import docx2txt

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────
# FastEmbed wrapper
# ─────────────────────────────────────────────

class FastEmbedWrapper(Embeddings):
    def __init__(self, model_name: str, cache_dir: Optional[str] = None):
        kwargs = {"cache_dir": cache_dir} if cache_dir else {}
        self._model = TextEmbedding(model_name, **kwargs)
        self.model_name = model_name

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        return [list(v) for v in self._model.embed(texts)]

    def embed_query(self, text: str) -> List[float]:
        return list(list(self._model.embed([text]))[0])


# ─────────────────────────────────────────────
# Client config
# ─────────────────────────────────────────────

CONFIG_PATH = os.getenv("CLIENT_CONFIG", "config.yaml")

def load_client_config() -> dict:
    defaults = {
        "branding": {
            "company_name": "RAG Assistant",
            "tagline": "Ask questions across your documents",
            "accent_colour": "#185FA5",
            "assistant_name": "Document Assistant",
        },
        "suggested_questions": [
            "What does this document cover?",
            "Summarise the key points",
            "What are the main steps?",
            "List the key findings",
        ],
        "ui": {
            "show_response_time": True,
            "show_sources": True,
            "max_suggested_questions": 5,
            "welcome_message": "Ask questions about your loaded documents.",
        },
    }
    try:
        p = Path(CONFIG_PATH)
        if p.exists():
            with open(p) as f:
                user = yaml.safe_load(f) or {}
            for section, values in user.items():
                if section in defaults and isinstance(values, dict):
                    defaults[section].update(values)
                else:
                    defaults[section] = values
    except Exception as e:
        logger.warning(f"Could not load {CONFIG_PATH}: {e}")
    return defaults


# ─────────────────────────────────────────────
# API key persistence
# ─────────────────────────────────────────────

API_KEY_FILE = "/app/.weldai_api_key"

def load_saved_api_key() -> str:
    try:
        p = Path(API_KEY_FILE)
        if p.exists():
            key = p.read_text().strip()
            if key.startswith("sk-ant-"):
                return key
    except Exception:
        pass
    return ""

def save_api_key(key: str):
    try:
        Path(API_KEY_FILE).write_text(key.strip())
    except Exception as e:
        logger.warning(f"Could not save API key: {e}")

def delete_api_key():
    try:
        p = Path(API_KEY_FILE)
        if p.exists():
            p.unlink()
    except Exception:
        pass


# ─────────────────────────────────────────────
# Claude pricing
# ─────────────────────────────────────────────

CLAUDE_PRICING: Dict[str, tuple] = {
    "claude-haiku-4-5-20251001": (0.80,   4.00),
    "claude-sonnet-4-6":         (3.00,  15.00),
    "claude-opus-4-6":          (15.00,  75.00),
}

def _model_rates(model: str):
    for key, (inp, out) in CLAUDE_PRICING.items():
        if model == key:
            return inp / 1_000_000, out / 1_000_000
    return 3.00 / 1_000_000, 15.00 / 1_000_000


# ─────────────────────────────────────────────
# System config
# ─────────────────────────────────────────────

@dataclass
class Config:
    docs_path: str         = field(default_factory=lambda: os.getenv("DOCS_PATH", "docs"))
    db_host: str           = field(default_factory=lambda: os.getenv("DB_HOST", "localhost"))
    db_port: str           = field(default_factory=lambda: os.getenv("DB_PORT", "5432"))
    db_name: str           = field(default_factory=lambda: os.getenv("DB_NAME", "ragdb"))
    db_user: str           = field(default_factory=lambda: os.getenv("DB_USER", "rag"))
    db_password: str       = field(default_factory=lambda: os.getenv("DB_PASSWORD", "ragpass"))
    collection_name: str   = field(default_factory=lambda: os.getenv("COLLECTION_NAME", "rag_docs"))
    embedding_model: str   = field(default_factory=lambda: os.getenv("EMBEDDING_MODEL", "BAAI/bge-small-en-v1.5"))
    chunk_size: int        = field(default_factory=lambda: int(os.getenv("CHUNK_SIZE", "1000")))
    chunk_overlap: int     = field(default_factory=lambda: int(os.getenv("CHUNK_OVERLAP", "100")))
    retrieval_k: int       = field(default_factory=lambda: int(os.getenv("RETRIEVAL_K", "8")))
    anthropic_api_key: str = field(default_factory=lambda: os.getenv("ANTHROPIC_API_KEY", "") or load_saved_api_key())
    claude_model: str      = field(default_factory=lambda: os.getenv("CLAUDE_MODEL", "claude-haiku-4-5-20251001"))
    max_tokens: int        = field(default_factory=lambda: int(os.getenv("MAX_TOKENS", "1024")))
    temperature: float     = field(default_factory=lambda: float(os.getenv("TEMPERATURE", "0.1")))
    max_question_len: int  = field(default_factory=lambda: int(os.getenv("MAX_QUESTION_LEN", "2000")))
    cache_ttl: int         = field(default_factory=lambda: int(os.getenv("CACHE_TTL", "86400")))

    @property
    def connection_string(self) -> str:
        return (
            f"postgresql+psycopg://{self.db_user}:{self.db_password}"
            f"@{self.db_host}:{self.db_port}/{self.db_name}"
        )

    def validate(self) -> List[str]:
        errors = []
        if not self.anthropic_api_key:
            errors.append("ANTHROPIC_API_KEY is not set")
        elif not self.anthropic_api_key.startswith("sk-ant-"):
            errors.append("ANTHROPIC_API_KEY doesn't look valid")
        return errors


# ─────────────────────────────────────────────
# Document loaders
# ─────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".txt", ".pdf", ".docx", ".pptx", ".xlsx", ".csv",
    ".doc", ".ppt", ".xls", ".rtf",
}

FILE_ICONS = {
    "pdf": "📄", "docx": "📝", "doc": "📝",
    "pptx": "📊", "ppt": "📊",
    "xlsx": "📋", "xls": "📋", "csv": "📋",
    "txt": "📃", "rtf": "📃",
}

def file_icon(name: str) -> str:
    return FILE_ICONS.get(Path(name).suffix.lower().lstrip("."), "📄")


def load_document(filepath: Path) -> List[Document]:
    ext = filepath.suffix.lower()
    if ext in {".doc", ".ppt", ".xls", ".rtf"}:
        converted = _convert_legacy_office(filepath)
        if converted:
            return load_document(converted)
        return []
    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(str(filepath))
            docs = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    docs.append(Document(
                        page_content=text,
                        metadata={"source": str(filepath), "page": i},
                    ))
            return docs
        elif ext == ".txt":
            text = filepath.read_text(encoding="utf-8", errors="replace")
            return [Document(page_content=text, metadata={"source": str(filepath)})]
        elif ext == ".docx":
            text = docx2txt.process(str(filepath))
            return [Document(page_content=text or "", metadata={"source": str(filepath)})]
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(filepath))
            text = [
                shape.text.strip()
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            return [Document(page_content=" ".join(text), metadata={"source": str(filepath)})]
        elif ext in (".xlsx", ".csv"):
            import pandas as pd
            df = pd.read_excel(filepath) if ext == ".xlsx" else pd.read_csv(filepath)
            return [Document(page_content=df.to_string(index=False), metadata={"source": str(filepath)})]
        else:
            return []
    except Exception as e:
        logger.error(f"Failed to load {filepath}: {e}")
        return []


def _convert_legacy_office(filepath: Path) -> Optional[Path]:
    import subprocess
    ext = filepath.suffix.lower()
    fmt = "docx" if ext in {".doc", ".rtf"} else ("pptx" if ext == ".ppt" else "xlsx")
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", fmt,
             "--outdir", str(filepath.parent), str(filepath)],
            capture_output=True, timeout=60,
        )
        if result.returncode == 0:
            converted = filepath.with_suffix("." + fmt)
            if converted.exists():
                return converted
    except Exception as e:
        logger.warning(f"LibreOffice conversion failed: {e}")
    return None


def _is_toc_chunk(text: str, page: int = 999) -> bool:
    """Detect table of contents chunks.
    Two signals:
    1. Chunk is from an early page (0-2) AND contains dotted leader lines (.... 31)
    2. More than 40% of lines match TOC pattern regardless of page
    """
    import re
    # TOC leader pattern: text followed by dots then page number
    toc_pattern = re.compile(r'.{5,}[.\s]{3,}\d+\s*$')
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if len(lines) < 2:
        return False
    toc_lines = sum(1 for l in lines if toc_pattern.match(l))
    ratio = toc_lines / len(lines)
    # Early pages with any TOC lines — skip
    if page <= 2 and toc_lines >= 2:
        return True
    # Any page with majority TOC lines — skip
    return ratio > 0.4


def collect_files(docs_path: str) -> List[Path]:
    root = Path(docs_path)
    if not root.exists():
        return []
    files = []
    for ext in SUPPORTED_EXTENSIONS:
        for f in root.rglob(f"*{ext}"):
            if any(part.startswith(".") for part in f.parts[1:]):
                continue
            files.append(f)
    return sorted(files)


# ─────────────────────────────────────────────
# Prompts
# ─────────────────────────────────────────────

SYSTEM_PROMPT = """\
You are a helpful AI assistant for business document search.
Answer questions using ONLY the context provided.
If the answer is not in the context, say: "I don't have enough information to answer that."

Rules:
- Read ALL context chunks before answering
- Be concise and specific
- Cite the source document name and page number where possible
- Do not make up or infer information not in the context"""

USER_PROMPT_TEMPLATE = "Context:\n{context}\n\nQuestion: {question}"


# ─────────────────────────────────────────────
# Metrics
# ─────────────────────────────────────────────

@dataclass
class APIMetrics:
    total_queries: int = 0
    input_tokens: int  = 0
    output_tokens: int = 0
    model: str = ""

    @property
    def total_tokens(self) -> int:
        return self.input_tokens + self.output_tokens

    @property
    def total_cost(self) -> float:
        in_rate, out_rate = _model_rates(self.model)
        return self.input_tokens * in_rate + self.output_tokens * out_rate

    def add_usage(self, inp: int, out: int):
        self.total_queries += 1
        self.input_tokens  += inp
        self.output_tokens += out


# ─────────────────────────────────────────────
# RAG System
# ─────────────────────────────────────────────

class RAGSystem:
    def __init__(self, cfg: Optional[Config] = None):
        self.cfg = cfg or Config()
        self.embeddings: Optional[FastEmbedWrapper] = None
        self.vectorstore: Optional[PGVector] = None
        self.retriever = None
        self._claude: Optional[anthropic.Anthropic] = None
        self._pool: Optional[pg_pool.SimpleConnectionPool] = None
        self.metrics = APIMetrics(model=self.cfg.claude_model)
        self._query_cache: Dict[str, Dict] = {}
        self.last_indexed: Optional[datetime] = None
        self._initialized = False

    # ── Embeddings ────────────────────────────────────────────────────

    def init_embeddings(self):
        if self.embeddings is None:
            logger.info(f"Loading embedding model: {self.cfg.embedding_model}")
            cache_dir = os.getenv("FASTEMBED_CACHE_PATH")
            self.embeddings = FastEmbedWrapper(self.cfg.embedding_model, cache_dir=cache_dir)

    # ── DB pool ───────────────────────────────────────────────────────

    def _init_pool(self):
        if self._pool is None:
            self._pool = pg_pool.SimpleConnectionPool(
                minconn=1, maxconn=5,
                host=self.cfg.db_host, port=self.cfg.db_port,
                dbname=self.cfg.db_name, user=self.cfg.db_user,
                password=self.cfg.db_password,
            )

    @contextlib.contextmanager
    def _db(self):
        self._init_pool()
        conn = self._pool.getconn()
        try:
            yield conn
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            self._pool.putconn(conn)

    # ── Vectorstore ───────────────────────────────────────────────────

    def connect_vectorstore(self):
        self.init_embeddings()
        self._init_pool()
        self.vectorstore = PGVector(
            embeddings=self.embeddings,
            collection_name=self.cfg.collection_name,
            connection=self.cfg.connection_string,
        )

    # ── Setup ─────────────────────────────────────────────────────────

    def setup(self, rebuild=False, progress_cb: Optional[Callable] = None) -> bool:
        try:
            _cb(progress_cb, 0.05, "Loading embedding model...")
            self.init_embeddings()
            _cb(progress_cb, 0.15, "Connecting to vector store...")
            self.connect_vectorstore()
            if rebuild == "incremental":
                _cb(progress_cb, 0.20, "Checking for changes...")
                self.index_documents(progress_cb=progress_cb, progress_offset=0.20, incremental=True)
            elif rebuild:
                _cb(progress_cb, 0.20, "Re-indexing all documents...")
                self.index_documents(progress_cb=progress_cb, progress_offset=0.20, incremental=False)
            else:
                try:
                    test = self.vectorstore.similarity_search("test", k=1)
                    if not test:
                        self.index_documents(progress_cb=progress_cb, progress_offset=0.20)
                except Exception:
                    self.index_documents(progress_cb=progress_cb, progress_offset=0.20)
            _cb(progress_cb, 0.85, f"Connecting to Claude AI...")
            self.build_qa_chain()
            _cb(progress_cb, 1.00, "Ready!")
            return True
        except Exception as e:
            logger.exception("RAG setup failed")
            _cb(progress_cb, 1.00, f"Failed: {e}")
            return False

    # ── Claude client ─────────────────────────────────────────────────

    def build_qa_chain(self):
        fresh_key = os.getenv("ANTHROPIC_API_KEY", "") or load_saved_api_key()
        if fresh_key:
            self.cfg.anthropic_api_key = fresh_key
        errors = self.cfg.validate()
        if errors:
            raise ValueError(f"Configuration errors: {'; '.join(errors)}")
        self._claude = anthropic.Anthropic(api_key=self.cfg.anthropic_api_key, max_retries=3)
        try:
            self._claude.models.list()
        except anthropic.AuthenticationError:
            raise ValueError("ANTHROPIC_API_KEY is invalid.")
        retrieve_k = self.cfg.retrieval_k
        # MMR with high lambda — mostly relevance, slight diversity to avoid TOC clustering
        self.retriever = self.vectorstore.as_retriever(
            search_type="mmr",
            search_kwargs={
                "k": retrieve_k,
                "fetch_k": retrieve_k * 4,  # fetch more candidates
                "lambda_mult": 0.9,          # 0=max diversity, 1=max relevance — keep high
            },
        )
        self._initialized = True

    # ── File hashing ──────────────────────────────────────────────────

    def _get_file_hash(self, filepath: Path) -> str:
        h = hashlib.md5()
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()

    def _load_hash_store(self) -> Dict[str, str]:
        hash_file = Path(self.cfg.docs_path) / ".file_hashes.json"
        try:
            if hash_file.exists():
                return json.loads(hash_file.read_text())
        except Exception:
            pass
        return {}

    def _save_hash_store(self, hashes: Dict[str, str]):
        try:
            hash_file = Path(self.cfg.docs_path) / ".file_hashes.json"
            hash_file.write_text(json.dumps(hashes, indent=2))
        except Exception as e:
            logger.warning(f"Could not save hash store: {e}")

    # ── Document indexing ─────────────────────────────────────────────

    def index_documents(
        self,
        progress_cb: Optional[Callable] = None,
        progress_offset: float = 0.0,
        incremental: bool = False,
    ) -> int:
        files = collect_files(self.cfg.docs_path)
        if not files:
            raise FileNotFoundError(f"No documents found in '{self.cfg.docs_path}'")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.cfg.chunk_size,
            chunk_overlap=self.cfg.chunk_overlap,
        )

        if incremental:
            stored_hashes = self._load_hash_store()
            current_hashes: Dict[str, str] = {}
            new_or_changed: List[Path] = []
            for fp in files:
                fh = self._get_file_hash(fp)
                current_hashes[str(fp)] = fh
                if stored_hashes.get(str(fp)) != fh:
                    new_or_changed.append(fp)
            deleted = [k for k in stored_hashes if not Path(k).exists()]
            if not new_or_changed and not deleted:
                _cb(progress_cb, progress_offset + 0.60, "All files up to date")
                self._persist_timestamp()
                return 0
            if files_to_remove := [str(fp) for fp in new_or_changed] + deleted:
                try:
                    with self._db() as conn:
                        cur = conn.cursor()
                        for source in files_to_remove:
                            cur.execute(
                                "DELETE FROM langchain_pg_embedding e "
                                "USING langchain_pg_collection c "
                                "WHERE e.collection_id = c.uuid AND c.name = %s "
                                "AND e.cmetadata->>'source' = %s",
                                (self.cfg.collection_name, source),
                            )
                except Exception as e:
                    logger.warning(f"Could not remove old chunks: {e}")
            files_to_index = new_or_changed
            self._save_hash_store(current_hashes)
        else:
            files_to_index = files
            current_hashes = {str(fp): self._get_file_hash(fp) for fp in files}

        all_chunks: List[Document] = []
        for i, fp in enumerate(files_to_index):
            frac = progress_offset + (0.55 * (i / max(len(files_to_index), 1)))
            _cb(progress_cb, frac, f"Reading {fp.name} ({i+1}/{len(files_to_index)})...")
            docs = load_document(fp)
            chunks = splitter.split_documents(docs)
            filtered = []
            skipped_toc = 0
            for chunk in chunks:
                page_num = chunk.metadata.get("page", 999)
                if _is_toc_chunk(chunk.page_content, page=int(page_num) if page_num != "" else 999):
                    skipped_toc += 1
                    continue
                source = Path(chunk.metadata.get("source", fp.name)).name
                page = chunk.metadata.get("page", "")
                page_str = f" | Page {int(page)+1}" if page != "" else ""
                chunk.page_content = f"[{source}{page_str}] " + chunk.page_content
                filtered.append(chunk)
            all_chunks.extend(filtered)
            logger.info(f"  {fp.name}: {len(docs)} pages → {len(filtered)} chunks (skipped {skipped_toc} TOC chunks)")

        _cb(progress_cb, progress_offset + 0.60, f"Embedding {len(all_chunks)} chunks...")

        EMBED_BATCH = 64  # increased from 32 — ~2x faster on 8GB+ RAM
        if not incremental:
            try:
                PGVector(
                    embeddings=self.embeddings,
                    collection_name=self.cfg.collection_name,
                    connection=self.cfg.connection_string,
                    pre_delete_collection=True,
                )
            except Exception:
                pass
            total_batches = (len(all_chunks) - 1) // EMBED_BATCH + 1
            for batch_num, i in enumerate(range(0, len(all_chunks), EMBED_BATCH)):
                batch = all_chunks[i:i + EMBED_BATCH]
                batch_progress = batch_num / max(total_batches, 1)
                _cb(progress_cb,
                    progress_offset + 0.60 + (0.22 * batch_progress),
                    f"Embedding batch {batch_num + 1}/{total_batches}...")
                self.vectorstore.add_documents(batch)
            self._save_hash_store(current_hashes)
        else:
            if all_chunks:
                for i in range(0, len(all_chunks), EMBED_BATCH):
                    self.vectorstore.add_documents(all_chunks[i:i + EMBED_BATCH])

        self._persist_timestamp()
        logger.info(f"Indexed {len(all_chunks)} chunks from {len(files_to_index)} files")
        return len(all_chunks)

    def _persist_timestamp(self):
        self.last_indexed = datetime.now()

    # ── Cache ─────────────────────────────────────────────────────────

    def _cache_key(self, q: str) -> str:
        return hashlib.md5(q.strip().lower().encode()).hexdigest()

    def _cache_get(self, q: str) -> Optional[Dict]:
        if self.cfg.cache_ttl <= 0:
            return None
        entry = self._query_cache.get(self._cache_key(q))
        if entry and (time.time() - entry["ts"]) < self.cfg.cache_ttl:
            return entry["result"]
        return None

    def _cache_set(self, q: str, result: Dict):
        if self.cfg.cache_ttl > 0:
            self._query_cache[self._cache_key(q)] = {"result": result, "ts": time.time()}

    # ── Retrieve ──────────────────────────────────────────────────────

    def _retrieve(self, question: str) -> List[Document]:
        return self.retriever.invoke(question)

    def _build_sources(self, docs: List[Document]) -> List[Dict]:
        sources = []
        for doc in docs:
            page = doc.metadata.get("page", "")
            name = Path(doc.metadata.get("source", "unknown")).name
            sources.append({
                "content": doc.page_content[:300],
                "metadata": doc.metadata,
                "display_name": f"{name} p.{int(page)+1}" if page != "" else name,
                "icon": file_icon(name),
            })
        return sources

    # ── Ask (streaming generator) ─────────────────────────────────────

    def ask_stream(self, question: str, history: Optional[List[Dict]] = None) -> Generator[Dict, None, None]:
        """Yield SSE-compatible dicts: {type, data}
        history: list of {role, content} dicts from previous exchanges (last N turns)
        """
        if not self._initialized:
            yield {"type": "error", "data": "System not initialised."}
            return

        question = question.strip()[:self.cfg.max_question_len]

        cached = self._cache_get(question)
        if cached:
            yield {"type": "cached", "data": cached["answer"]}
            yield {"type": "sources", "data": cached["sources"]}
            yield {"type": "meta", "data": {"elapsed": cached["elapsed"], "cached": True}}
            yield {"type": "done", "data": ""}
            return

        try:
            t0 = time.time()
            yield {"type": "status", "data": "Searching documents..."}
            source_docs = self._retrieve(question)
            yield {"type": "status", "data": "Generating answer..."}
            context = "\n\n".join(d.page_content for d in source_docs)

            # Build messages array from history
            # Frontend sends history including the current question as last item
            # We replace the last user message with the RAG-augmented version
            messages = []
            if history and len(history) >= 1:
                # Add all history except the last user message (we'll add RAG version)
                prior = [h for h in history[:-1] if h.get("role") in ("user", "assistant") and h.get("content")]
                for turn in prior[-10:]:
                    messages.append({"role": turn["role"], "content": turn["content"]})
            # Always end with the RAG-augmented current question
            messages.append({"role": "user", "content": USER_PROMPT_TEMPLATE.format(
                context=context, question=question
            )})

            answer_buf = ""
            with self._claude.messages.stream(
                model=self.cfg.claude_model,
                max_tokens=self.cfg.max_tokens,
                temperature=self.cfg.temperature,
                system=SYSTEM_PROMPT,
                messages=messages,
            ) as stream:
                for text in stream.text_stream:
                    answer_buf += text
                    yield {"type": "token", "data": text}
                final = stream.get_final_message()
                self.metrics.add_usage(final.usage.input_tokens, final.usage.output_tokens)

            elapsed = round(time.time() - t0, 1)
            sources = self._build_sources(source_docs)
            result = {"answer": answer_buf, "sources": sources, "elapsed": elapsed}
            self._cache_set(question, result)

            yield {"type": "sources", "data": sources}
            yield {"type": "meta", "data": {
                "elapsed": elapsed,
                "cached": False,
                "cost": round(self.metrics.total_cost, 5),
                "tokens": self.metrics.total_tokens,
            }}
            yield {"type": "done", "data": ""}

        except anthropic.APIConnectionError:
            yield {"type": "error", "data": "Could not reach the Anthropic API. Check your internet connection."}
        except anthropic.RateLimitError:
            yield {"type": "error", "data": "Rate limit reached. Please wait a moment and try again."}
        except anthropic.AuthenticationError:
            yield {"type": "error", "data": "Invalid API key. Check your configuration."}
        except Exception as e:
            logger.error(f"Query failed: {e}")
            yield {"type": "error", "data": str(e)}

    # ── Stats ─────────────────────────────────────────────────────────

    def get_stats(self) -> Dict[str, Any]:
        count = 0
        try:
            with self._db() as conn:
                cur = conn.cursor()
                cur.execute(
                    "SELECT COUNT(*) FROM langchain_pg_embedding e "
                    "JOIN langchain_pg_collection c ON e.collection_id = c.uuid "
                    "WHERE c.name = %s",
                    (self.cfg.collection_name,),
                )
                count = cur.fetchone()[0]
        except Exception:
            pass
        files = collect_files(self.cfg.docs_path)
        m = self.metrics
        return {
            "chunks_indexed": count,
            "documents_found": len(files),
            "embedding_model": self.cfg.embedding_model,
            "llm_model": self.cfg.claude_model,
            "chunk_size": self.cfg.chunk_size,
            "retrieval_k": self.cfg.retrieval_k,
            "cache_size": len(self._query_cache),
            "last_indexed": (
                self.last_indexed.strftime("%Y-%m-%d %H:%M") if self.last_indexed else "Never"
            ),
            "session_queries": m.total_queries,
            "session_tokens": m.total_tokens,
            "session_cost_usd": round(m.total_cost, 5),
        }

    def get_changes(self) -> Dict[str, Any]:
        files = collect_files(self.cfg.docs_path)
        stored_hashes = self._load_hash_store()
        if not stored_hashes:
            return {"new": [], "deleted": [], "has_changes": False}
        stored_by_name = {Path(k).name: v for k, v in stored_hashes.items()}
        new_or_changed = [f.name for f in files if stored_by_name.get(f.name) != self._get_file_hash(f)]
        deleted = list(set(Path(k).name for k in stored_hashes) - set(f.name for f in files))
        return {
            "new": new_or_changed,
            "deleted": deleted,
            "has_changes": bool(new_or_changed or deleted),
        }

    def get_chunks_in_db(self) -> int:
        try:
            conn = psycopg2.connect(
                host=self.cfg.db_host, port=int(self.cfg.db_port),
                dbname=self.cfg.db_name, user=self.cfg.db_user,
                password=self.cfg.db_password,
            )
            cur = conn.cursor()
            cur.execute(
                "SELECT COUNT(*) FROM langchain_pg_embedding e "
                "JOIN langchain_pg_collection c ON e.collection_id = c.uuid "
                "WHERE c.name = %s",
                (self.cfg.collection_name,),
            )
            count = cur.fetchone()[0]
            conn.close()
            return count
        except Exception:
            return 0

    def clear_cache(self):
        self._query_cache.clear()

    # ── Chat history ──────────────────────────────────────────────────

    def ensure_history_table(self):
        """Create chat_messages table if it doesn't exist."""
        try:
            with self._db() as conn:
                cur = conn.cursor()
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS chat_messages (
                        id SERIAL PRIMARY KEY,
                        session_id VARCHAR(64) NOT NULL,
                        role VARCHAR(16) NOT NULL,
                        content TEXT NOT NULL,
                        sources JSONB,
                        elapsed FLOAT,
                        cached BOOLEAN DEFAULT FALSE,
                        created_at TIMESTAMPTZ DEFAULT NOW()
                    )
                """)
                cur.execute("""
                    CREATE INDEX IF NOT EXISTS idx_chat_messages_session
                    ON chat_messages(session_id, created_at)
                """)
        except Exception as e:
            logger.warning(f"Could not create chat_messages table: {e}")

    def save_message(self, session_id: str, role: str, content: str,
                     sources: Optional[List] = None, elapsed: Optional[float] = None,
                     cached: bool = False):
        try:
            with self._db() as conn:
                cur = conn.cursor()
                cur.execute(
                    "INSERT INTO chat_messages (session_id, role, content, sources, elapsed, cached) "
                    "VALUES (%s, %s, %s, %s, %s, %s)",
                    (session_id, role, content,
                     json.dumps(sources) if sources else None,
                     elapsed, cached)
                )
        except Exception as e:
            logger.warning(f"Could not save message: {e}")

    def get_history(self, session_id: str, limit: int = 50) -> List[Dict]:
        try:
            with self._db() as conn:
                cur = conn.cursor()
                cur.execute(
                    "SELECT role, content, sources, elapsed, cached, created_at "
                    "FROM chat_messages WHERE session_id = %s "
                    "ORDER BY created_at ASC LIMIT %s",
                    (session_id, limit)
                )
                rows = cur.fetchall()
                return [
                    {
                        "role": r[0],
                        "content": r[1],
                        "sources": json.loads(r[2]) if r[2] else [],
                        "elapsed": r[3],
                        "cached": r[4],
                        "created_at": r[5].isoformat() if r[5] else None,
                    }
                    for r in rows
                ]
        except Exception as e:
            logger.warning(f"Could not get history: {e}")
            return []

    def clear_history(self, session_id: str):
        try:
            with self._db() as conn:
                cur = conn.cursor()
                cur.execute("DELETE FROM chat_messages WHERE session_id = %s", (session_id,))
        except Exception as e:
            logger.warning(f"Could not clear history: {e}")

    def list_sessions(self) -> List[Dict]:
        try:
            with self._db() as conn:
                cur = conn.cursor()
                cur.execute("""
                    SELECT session_id,
                           COUNT(*) FILTER (WHERE role='user') as questions,
                           MIN(created_at) as started,
                           MAX(created_at) as last_active
                    FROM chat_messages
                    GROUP BY session_id
                    ORDER BY last_active DESC
                    LIMIT 20
                """)
                rows = cur.fetchall()
                return [
                    {
                        "session_id": r[0],
                        "questions": r[1],
                        "started": r[2].isoformat() if r[2] else None,
                        "last_active": r[3].isoformat() if r[3] else None,
                    }
                    for r in rows
                ]
        except Exception as e:
            logger.warning(f"Could not list sessions: {e}")
            return []


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def _cb(cb: Optional[Callable], progress: float, message: str):
    if cb:
        cb(progress, message)
    logger.info(f"[{progress:.0%}] {message}")