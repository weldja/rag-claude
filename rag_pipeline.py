"""
RAG Pipeline — Claude API Edition  (Production)
────────────────────────────────────────────────
Vectorstore : PostgreSQL + pgvector
Embeddings  : BAAI/bge-small-en-v1.5  via fastembed  (local, no API cost)
LLM         : Anthropic Claude API    (real streaming + token tracking)
"""

import contextlib
import os
import threading
import sys
import time
import hashlib
import json
import logging
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, field

import streamlit as st
import yaml
import anthropic
import psycopg2
import psycopg2.pool as pg_pool

# LangChain — chunking + vectorstore only (no community package)
from langchain_postgres import PGVector
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings

# Direct library imports (replaces langchain-community loaders)
import pypdf
import docx2txt
from fastembed import TextEmbedding

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────
# FastEmbed wrapper (implements LangChain Embeddings interface)
# Replaces langchain_community FastEmbedEmbeddings — no deprecation warnings
# ─────────────────────────────────────────────

class FastEmbedWrapper(Embeddings):
    """Thin wrapper around fastembed.TextEmbedding for LangChain compatibility."""

    def __init__(self, model_name: str, cache_dir: Optional[str] = None):
        kwargs = {}
        if cache_dir:
            kwargs["cache_dir"] = cache_dir
        self._model = TextEmbedding(model_name, **kwargs)
        self.model_name = model_name

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        return [list(v) for v in self._model.embed(texts)]

    def embed_query(self, text: str) -> List[float]:
        return list(list(self._model.embed([text]))[0])

# ─────────────────────────────────────────────
# Client config loader
# ─────────────────────────────────────────────

CONFIG_PATH = os.getenv("CLIENT_CONFIG", "config.yaml")

def load_client_config() -> dict:
    defaults = {
        "branding": {
            "company_name": "RAG Assistant",
            "tagline": "Ask questions across your documents",
            "accent_colour": "#185FA5",
            "assistant_name": "Document Assistant",
        },
        "suggested_questions": [
            "What does this document cover?",
            "Summarise the key points",
            "What are the main steps?",
        ],
        "ui": {
            "show_response_time": True,
            "show_sources": True,
            "show_copy_button": True,
            "max_suggested_questions": 5,
            "welcome_message": "Ask questions about your loaded documents.",
        },
    }
    try:
        if Path(CONFIG_PATH).exists():
            with open(CONFIG_PATH) as f:
                user = yaml.safe_load(f) or {}
            for section, values in user.items():
                if section in defaults and isinstance(values, dict):
                    defaults[section].update(values)
                else:
                    defaults[section] = values
    except Exception as e:
        logger.warning(f"Could not load {CONFIG_PATH}: {e}")
    return defaults



# ─────────────────────────────────────────────
# API key persistence (browser-entered key)
# ─────────────────────────────────────────────

API_KEY_FILE = "/app/.weldai_api_key"

def load_saved_api_key() -> str:
    """Load API key saved via the browser UI."""
    try:
        p = Path(API_KEY_FILE)
        if p.exists():
            key = p.read_text().strip()
            if key.startswith("sk-ant-"):
                return key
    except Exception:
        pass
    return ""

def save_api_key(key: str):
    """Persist API key entered via browser to a local file."""
    try:
        Path(API_KEY_FILE).write_text(key.strip())
    except Exception as e:
        logger.warning(f"Could not save API key: {e}")

def delete_api_key():
    """Remove saved API key."""
    try:
        p = Path(API_KEY_FILE)
        if p.exists():
            p.unlink()
    except Exception:
        pass

# ─────────────────────────────────────────────
# Claude pricing  ($ per million tokens)
# ─────────────────────────────────────────────

CLAUDE_PRICING: Dict[str, tuple] = {
    "claude-haiku-4-5-20251001": (0.80,   4.00),
    "claude-sonnet-4-6":         (3.00,  15.00),
    "claude-opus-4-6":          (15.00,  75.00),
}

def _model_rates(model: str):
    for key, (inp, out) in CLAUDE_PRICING.items():
        if model == key:
            return inp / 1_000_000, out / 1_000_000
    return 3.00 / 1_000_000, 15.00 / 1_000_000  # safe default


# ─────────────────────────────────────────────
# System configuration
# ─────────────────────────────────────────────

@dataclass
class Config:
    docs_path: str         = field(default_factory=lambda: os.getenv("DOCS_PATH", "docs"))
    db_host: str           = field(default_factory=lambda: os.getenv("DB_HOST", "localhost"))
    db_port: str           = field(default_factory=lambda: os.getenv("DB_PORT", "5432"))
    db_name: str           = field(default_factory=lambda: os.getenv("DB_NAME", "ragdb"))
    db_user: str           = field(default_factory=lambda: os.getenv("DB_USER", "rag"))
    db_password: str       = field(default_factory=lambda: os.getenv("DB_PASSWORD", "ragpass"))
    collection_name: str   = field(default_factory=lambda: os.getenv("COLLECTION_NAME", "rag_docs"))
    embedding_model: str   = field(default_factory=lambda: os.getenv("EMBEDDING_MODEL", "BAAI/bge-small-en-v1.5"))
    chunk_size: int        = field(default_factory=lambda: int(os.getenv("CHUNK_SIZE", "1000")))
    chunk_overlap: int     = field(default_factory=lambda: int(os.getenv("CHUNK_OVERLAP", "100")))
    retrieval_k: int       = field(default_factory=lambda: int(os.getenv("RETRIEVAL_K", "4")))
    reranker_model: str    = field(default_factory=lambda: os.getenv("RERANKER_MODEL", ""))
    reranker_top_n: int    = field(default_factory=lambda: int(os.getenv("RERANKER_TOP_N", "4")))
    anthropic_api_key: str = field(default_factory=lambda: os.getenv("ANTHROPIC_API_KEY", "") or load_saved_api_key())
    claude_model: str      = field(default_factory=lambda: os.getenv("CLAUDE_MODEL", "claude-haiku-4-5-20251001"))
    max_tokens: int        = field(default_factory=lambda: int(os.getenv("MAX_TOKENS", "1024")))
    temperature: float     = field(default_factory=lambda: float(os.getenv("TEMPERATURE", "0.1")))
    max_question_len: int  = field(default_factory=lambda: int(os.getenv("MAX_QUESTION_LEN", "2000")))
    batch_size: int        = field(default_factory=lambda: int(os.getenv("BATCH_SIZE", "100")))
    cache_ttl: int         = field(default_factory=lambda: int(os.getenv("CACHE_TTL", "86400")))

    @property
    def connection_string(self) -> str:
        return (
            f"postgresql+psycopg://{self.db_user}:{self.db_password}"
            f"@{self.db_host}:{self.db_port}/{self.db_name}"
        )

    def validate(self) -> List[str]:
        """Return list of configuration errors (empty = valid)."""
        errors = []
        if not self.anthropic_api_key:
            errors.append("ANTHROPIC_API_KEY is not set")
        elif not self.anthropic_api_key.startswith("sk-ant-"):
            errors.append("ANTHROPIC_API_KEY doesn't look valid (should start with sk-ant-)")
        if self.chunk_size < 100:
            errors.append("CHUNK_SIZE should be at least 100")
        if self.retrieval_k < 1:
            errors.append("RETRIEVAL_K should be at least 1")
        return errors


config = Config()


# ─────────────────────────────────────────────
# Cost / usage tracking
# ─────────────────────────────────────────────

@dataclass
class APIMetrics:
    total_queries: int = 0
    input_tokens: int  = 0
    output_tokens: int = 0
    model: str = field(default_factory=lambda: config.claude_model)

    @property
    def total_tokens(self) -> int:
        return self.input_tokens + self.output_tokens

    @property
    def total_cost(self) -> float:
        in_rate, out_rate = _model_rates(self.model)
        return self.input_tokens * in_rate + self.output_tokens * out_rate

    def add_usage(self, input_tokens: int, output_tokens: int):
        self.total_queries += 1
        self.input_tokens  += input_tokens
        self.output_tokens += output_tokens


# ─────────────────────────────────────────────
# Prompts — system/user separated for Claude
# ─────────────────────────────────────────────

# System prompt: stable instructions → good for prompt caching in future
SYSTEM_PROMPT = """\
You are a helpful AI assistant for business document search.
Answer questions using ONLY the context provided.
If the answer is not in the context, say exactly: \
"I don't have enough information to answer that."

Rules:
- Read ALL context chunks before answering
- Prioritise chunks that directly answer the question
- Be concise and specific
- Cite the source document name and page number where possible
- Do not make up or infer information not present in the context
- If multiple chunks give conflicting answers, use the most relevant one"""

# User prompt: context + question only
USER_PROMPT_TEMPLATE = """\
Context:
{context}

Question: {question}"""


# ─────────────────────────────────────────────
# Document loaders
# ─────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".txt", ".pdf", ".docx", ".pptx", ".xlsx", ".csv",
    ".doc", ".ppt", ".xls", ".rtf",
}


def load_document(filepath: Path) -> List[Document]:
    ext = filepath.suffix.lower()
    if ext in {".doc", ".ppt", ".xls", ".rtf"}:
        converted = _convert_legacy_office(filepath)
        if converted:
            return load_document(converted)
        logger.warning(f"Could not convert {filepath.name} — skipping")
        return []
    try:
        if ext == ".pdf":
            reader = pypdf.PdfReader(str(filepath))
            docs = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    docs.append(Document(
                        page_content=text,
                        metadata={"source": str(filepath), "page": i},
                    ))
            return docs
        elif ext == ".txt":
            text = Path(filepath).read_text(encoding="utf-8", errors="replace")
            return [Document(page_content=text, metadata={"source": str(filepath)})]
        elif ext == ".docx":
            text = docx2txt.process(str(filepath))
            return [Document(page_content=text or "", metadata={"source": str(filepath)})]
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(filepath))
            text = [
                shape.text.strip()
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            return [Document(page_content=" ".join(text), metadata={"source": str(filepath)})]
        elif ext in (".xlsx", ".csv"):
            import pandas as pd
            df = pd.read_excel(filepath) if ext == ".xlsx" else pd.read_csv(filepath)
            return [Document(page_content=df.to_string(index=False), metadata={"source": str(filepath)})]
        else:
            return []
    except Exception as e:
        logger.error(f"Failed to load {filepath}: {e}")
        return []


def _convert_legacy_office(filepath: Path) -> Optional[Path]:
    import subprocess
    ext = filepath.suffix.lower()
    fmt = "docx" if ext in {".doc", ".rtf"} else ("pptx" if ext == ".ppt" else "xlsx")
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", fmt,
             "--outdir", str(filepath.parent), str(filepath)],
            capture_output=True, timeout=60,
        )
        if result.returncode == 0:
            converted = filepath.with_suffix("." + fmt)
            if converted.exists():
                return converted
    except Exception as e:
        logger.warning(f"LibreOffice conversion failed for {filepath}: {e}")
    return None


def collect_files(docs_path: str) -> List[Path]:
    root = Path(docs_path)
    if not root.exists():
        return []
    files = []
    for ext in SUPPORTED_EXTENSIONS:
        for f in root.rglob(f"*{ext}"):
            if any(part.startswith(".") for part in f.parts[1:]):
                continue
            files.append(f)
    return sorted(files)


# ─────────────────────────────────────────────
# RAG System
# ─────────────────────────────────────────────

class RAGSystem:
    def __init__(self, cfg: Config = config):
        self.cfg = cfg
        self.embeddings: Optional[FastEmbedWrapper] = None
        self.vectorstore: Optional[PGVector] = None
        self.retriever = None
        self.reranker = None
        self._claude: Optional[anthropic.Anthropic] = None
        self._pool: Optional[pg_pool.SimpleConnectionPool] = None
        self.metrics = APIMetrics(model=cfg.claude_model)
        self._query_cache: Dict[str, Dict] = {}
        self.last_indexed: Optional[datetime] = None
        self._initialized = False

    # ── Embeddings (fastembed — no PyTorch) ───────────────────────────

    def init_embeddings(self):
        if self.embeddings is None:
            logger.info(f"Loading embedding model via fastembed: {self.cfg.embedding_model}")
            cache_dir = os.getenv("FASTEMBED_CACHE_PATH", None)
            self.embeddings = FastEmbedWrapper(
                model_name=self.cfg.embedding_model,
                cache_dir=cache_dir,
            )

    # ── PostgreSQL connection pool ────────────────────────────────────

    def _init_pool(self):
        if self._pool is None:
            self._pool = pg_pool.SimpleConnectionPool(
                minconn=1,
                maxconn=5,
                host=self.cfg.db_host,
                port=self.cfg.db_port,
                dbname=self.cfg.db_name,
                user=self.cfg.db_user,
                password=self.cfg.db_password,
            )

    @contextlib.contextmanager
    def _db(self):
        """Yield a pooled database connection, return it when done."""
        self._init_pool()
        conn = self._pool.getconn()
        try:
            yield conn
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            self._pool.putconn(conn)

    # ── Vectorstore ───────────────────────────────────────────────────

    def connect_vectorstore(self):
        self.init_embeddings()
        self._init_pool()
        try:
            ts_file = Path(self.cfg.docs_path) / ".last_indexed"
            if ts_file.exists():
                self.last_indexed = datetime.fromisoformat(ts_file.read_text().strip())
        except Exception:
            pass
        self.vectorstore = PGVector(
            embeddings=self.embeddings,
            collection_name=self.cfg.collection_name,
            connection=self.cfg.connection_string,
        )

    # ── QA setup (retriever + Claude client) ─────────────────────────

    def build_qa_chain(self):
        if self.vectorstore is None:
            raise RuntimeError("Vectorstore not ready.")

        # Always read the freshest key — env var takes priority, then saved file
        fresh_key = os.getenv("ANTHROPIC_API_KEY", "") or load_saved_api_key()
        if fresh_key:
            self.cfg.anthropic_api_key = fresh_key

        errors = self.cfg.validate()
        if errors:
            raise ValueError(f"Configuration errors: {'; '.join(errors)}")

        # Anthropic client — max_retries handles transient network/overload errors
        self._claude = anthropic.Anthropic(
            api_key=self.cfg.anthropic_api_key,
            max_retries=3,
        )

        # Verify the key works (lightweight call, no cost)
        try:
            self._claude.models.list()
            logger.info(f"Claude API key verified. Model: {self.cfg.claude_model}")
        except anthropic.AuthenticationError:
            raise ValueError("ANTHROPIC_API_KEY is invalid. Check your .env file.")

        # Retriever — fetch more candidates when reranking
        retrieve_k = self.cfg.retrieval_k * 4 if self.cfg.reranker_model else self.cfg.retrieval_k
        self.retriever = self.vectorstore.as_retriever(
            search_type="similarity",
            search_kwargs={"k": retrieve_k},
        )

        # Optional cross-encoder reranker
        if self.cfg.reranker_model:
            try:
                from sentence_transformers import CrossEncoder
                model_path = Path("/app") / self.cfg.reranker_model
                if "models--" in str(model_path):
                    snapshots = model_path / "snapshots"
                    if snapshots.exists():
                        snap_dirs = list(snapshots.iterdir())
                        if snap_dirs:
                            model_path = snap_dirs[0]
                self.reranker = CrossEncoder(str(model_path))
                logger.info(f"Reranker loaded: {model_path}")
            except Exception as e:
                logger.error(f"Reranker failed to load: {e}")
                self.reranker = None

        self._initialized = True

    # ── Setup ─────────────────────────────────────────────────────────

    def setup(self, rebuild=False, progress_cb=None) -> bool:
        try:
            _cb(progress_cb, 0.05, "Loading embedding model...")
            self.init_embeddings()
            _cb(progress_cb, 0.15, "Connecting to vector store...")
            self.connect_vectorstore()
            if rebuild == "incremental":
                _cb(progress_cb, 0.20, "Checking for changes...")
                self.index_documents(progress_cb=progress_cb, progress_offset=0.20, incremental=True)
            elif rebuild:
                _cb(progress_cb, 0.20, "Re-indexing all documents...")
                self.index_documents(progress_cb=progress_cb, progress_offset=0.20, incremental=False)
            else:
                try:
                    test = self.vectorstore.similarity_search("test", k=1)
                    if not test:
                        _cb(progress_cb, 0.20, "No data found — indexing documents...")
                        self.index_documents(progress_cb=progress_cb, progress_offset=0.20)
                except Exception:
                    _cb(progress_cb, 0.20, "No data found — indexing documents...")
                    self.index_documents(progress_cb=progress_cb, progress_offset=0.20)
            _cb(progress_cb, 0.85, f"Connecting to Claude API ({self.cfg.claude_model})...")
            self.build_qa_chain()
            _cb(progress_cb, 1.00, "Ready!")
            return True
        except Exception as e:
            logger.exception("RAG setup failed")
            _cb(progress_cb, 1.00, f"Failed: {e}")
            return False

    # ── File hashing ──────────────────────────────────────────────────

    def _get_file_hash(self, filepath: Path) -> str:
        h = hashlib.md5()
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()

    def _load_hash_store(self) -> Dict[str, str]:
        hash_file = Path(self.cfg.docs_path) / ".file_hashes.json"
        try:
            if hash_file.exists():
                return json.loads(hash_file.read_text())
        except Exception:
            pass
        return {}

    def _save_hash_store(self, hashes: Dict[str, str]):
        hash_file = Path(self.cfg.docs_path) / ".file_hashes.json"
        try:
            hash_file.write_text(json.dumps(hashes, indent=2))
        except Exception as e:
            logger.warning(f"Could not save hash store: {e}")

    # ── Document indexing ─────────────────────────────────────────────

    def index_documents(
        self,
        progress_cb=None,
        progress_offset: float = 0.0,
        incremental: bool = False,
    ) -> int:
        files = collect_files(self.cfg.docs_path)
        if not files:
            raise FileNotFoundError(f"No supported documents found in '{self.cfg.docs_path}'")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.cfg.chunk_size,
            chunk_overlap=self.cfg.chunk_overlap,
            length_function=len,
        )

        if incremental:
            stored_hashes = self._load_hash_store()
            current_hashes: Dict[str, str] = {}
            new_or_changed: List[Path] = []

            for fp in files:
                fh = self._get_file_hash(fp)
                current_hashes[str(fp)] = fh
                if stored_hashes.get(str(fp)) != fh:
                    new_or_changed.append(fp)

            deleted = [k for k in stored_hashes if not Path(k).exists()]
            _cb(
                progress_cb, progress_offset + 0.05,
                f"Incremental: {len(new_or_changed)} changed, "
                f"{len(files) - len(new_or_changed)} unchanged, {len(deleted)} deleted",
            )

            if not new_or_changed and not deleted:
                _cb(progress_cb, progress_offset + 0.60, "All files up to date — nothing to index")
                self._persist_timestamp()
                return 0

            # Remove stale chunks from vectorstore using connection pool
            files_to_remove = [str(fp) for fp in new_or_changed] + deleted
            if files_to_remove:
                try:
                    with self._db() as conn:
                        cur = conn.cursor()
                        for source in files_to_remove:
                            cur.execute(
                                "DELETE FROM langchain_pg_embedding e "
                                "USING langchain_pg_collection c "
                                "WHERE e.collection_id = c.uuid "
                                "AND c.name = %s "
                                "AND e.cmetadata->>'source' = %s",
                                (self.cfg.collection_name, source),
                            )
                            logger.info(f"Removed chunks for: {source}")
                except Exception as e:
                    logger.warning(f"Could not remove old chunks: {e}")

            files_to_index = new_or_changed
            self._save_hash_store(current_hashes)
        else:
            files_to_index = files
            current_hashes = {str(fp): self._get_file_hash(fp) for fp in files}

        all_chunks: List[Document] = []
        for i, fp in enumerate(files_to_index):
            frac = (
                progress_offset + (0.60 * (i / len(files_to_index)))
                if files_to_index else progress_offset
            )
            _cb(progress_cb, frac, f"Loading {fp.name} ({i+1}/{len(files_to_index)})...")
            docs = load_document(fp)
            chunks = splitter.split_documents(docs)
            for chunk in chunks:
                source = Path(chunk.metadata.get("source", fp.name)).name
                page = chunk.metadata.get("page", "")
                page_str = f" | Page {int(page)+1}" if page != "" else ""
                chunk.page_content = f"[{source}{page_str}] " + chunk.page_content
            all_chunks.extend(chunks)
            logger.info(f"  {fp.name}: {len(docs)} pages → {len(chunks)} chunks")

        _cb(progress_cb, progress_offset + 0.60, f"Embedding {len(all_chunks)} chunks...")

        EMBED_BATCH = 32  # embed and insert in small batches to avoid memory spikes

        if not incremental:
            # Delete existing collection first
            try:
                PGVector(
                    embeddings=self.embeddings,
                    collection_name=self.cfg.collection_name,
                    connection=self.cfg.connection_string,
                    pre_delete_collection=True,
                )
            except Exception:
                pass
            # Insert in batches
            for i in range(0, len(all_chunks), EMBED_BATCH):
                batch = all_chunks[i:i + EMBED_BATCH]
                logger.info(f"  Embedding batch {i//EMBED_BATCH + 1}/{(len(all_chunks)-1)//EMBED_BATCH + 1} ({len(batch)} chunks)...")
                self.vectorstore.add_documents(batch)
            self._save_hash_store(current_hashes)
        else:
            if all_chunks:
                for i in range(0, len(all_chunks), EMBED_BATCH):
                    batch = all_chunks[i:i + EMBED_BATCH]
                    self.vectorstore.add_documents(batch)

        self._persist_timestamp()
        logger.info(f"Indexed {len(all_chunks)} chunks from {len(files_to_index)} files")
        return len(all_chunks)

    def _persist_timestamp(self):
        self.last_indexed = datetime.now()
        for path in ["/tmp/.last_indexed", "/app/docs/.last_indexed"]:
            try:
                with open(path, "w") as f:
                    f.write(self.last_indexed.isoformat())
                break
            except Exception:
                continue

    # ── Query cache ───────────────────────────────────────────────────

    def _cache_key(self, question: str) -> str:
        return hashlib.md5(question.strip().lower().encode()).hexdigest()

    def _cache_get(self, question: str) -> Optional[Dict]:
        if self.cfg.cache_ttl <= 0:
            return None
        key = self._cache_key(question)
        entry = self._query_cache.get(key)
        if entry and (time.time() - entry["ts"]) < self.cfg.cache_ttl:
            return entry["result"]
        return None

    def _cache_set(self, question: str, result: Dict):
        if self.cfg.cache_ttl > 0:
            self._query_cache[self._cache_key(question)] = {"result": result, "ts": time.time()}

    # ── Input sanitisation ────────────────────────────────────────────

    def _sanitise(self, question: str) -> str:
        q = question.strip()
        if len(q) > self.cfg.max_question_len:
            q = q[: self.cfg.max_question_len]
            logger.warning("Question truncated to max length")
        return q

    # ── Retrieve + rerank (shared by ask and ask_streaming) ───────────

    def _retrieve(self, question: str) -> List[Document]:
        source_docs = self.retriever.invoke(question)
        if self.reranker and source_docs:
            pairs = [(question, doc.page_content) for doc in source_docs]
            scores = self.reranker.predict(pairs)
            scored = sorted(zip(scores, source_docs), key=lambda x: x[0], reverse=True)
            source_docs = [doc for _, doc in scored[: self.cfg.reranker_top_n]]
            logger.info(f"Reranked → {len(source_docs)} chunks")
        return source_docs

    def _build_sources(self, source_docs: List[Document]) -> List[Dict]:
        sources = []
        for doc in source_docs:
            page = doc.metadata.get("page", "")
            name = Path(doc.metadata.get("source", "unknown")).name
            sources.append({
                "content": doc.page_content[:300],
                "metadata": doc.metadata,
                "display_name": f"{name} p.{int(page)+1}" if page != "" else name,
            })
        return sources

    # ── Ask (non-streaming — used by CLI and cache reads) ─────────────

    def ask(self, question: str) -> Dict[str, Any]:
        if not self._initialized:
            return {"answer": "System not initialised. Click Init in the sidebar.",
                    "sources": [], "cached": False, "error": True}

        question = self._sanitise(question)
        cached = self._cache_get(question)
        if cached:
            logger.info("Cache hit")
            return {**cached, "cached": True}

        try:
            t0 = time.time()
            source_docs = self._retrieve(question)
            context = "\n\n".join(d.page_content for d in source_docs)

            response = self._claude.messages.create(
                model=self.cfg.claude_model,
                max_tokens=self.cfg.max_tokens,
                temperature=self.cfg.temperature,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": USER_PROMPT_TEMPLATE.format(
                    context=context, question=question
                )}],
            )
            answer = response.content[0].text
            self.metrics.add_usage(response.usage.input_tokens, response.usage.output_tokens)
            logger.info(
                f"Claude — in:{response.usage.input_tokens} out:{response.usage.output_tokens} "
                f"session_cost:${self.metrics.total_cost:.5f}"
            )

            result = {
                "answer": answer,
                "sources": self._build_sources(source_docs),
                "elapsed": round(time.time() - t0, 1),
                "cached": False,
                "error": False,
            }
            self._cache_set(question, result)
            return result

        except anthropic.APIConnectionError:
            return {"answer": "⚠️ Could not reach the Anthropic API. Check your internet connection.",
                    "sources": [], "cached": False, "error": True}
        except anthropic.RateLimitError:
            return {"answer": "⚠️ Rate limit reached. Please wait a moment and try again.",
                    "sources": [], "cached": False, "error": True}
        except anthropic.AuthenticationError:
            return {"answer": "⚠️ Invalid API key. Check ANTHROPIC_API_KEY in your .env file.",
                    "sources": [], "cached": False, "error": True}
        except Exception as e:
            logger.error(f"Query failed: {e}")
            return {"answer": f"Error: {e}", "sources": [], "cached": False, "error": True}

    # ── Ask (streaming — used by the UI for live token output) ────────

    def ask_streaming(self, question: str, placeholder) -> Dict[str, Any]:
        """Stream Claude's response token-by-token into a Streamlit placeholder."""
        if not self._initialized:
            placeholder.error("System not initialised. Click Init in the sidebar.")
            return {"answer": "", "sources": [], "elapsed": 0, "cached": False, "error": True}

        question = self._sanitise(question)

        # Return cached result (fake-stream for consistency)
        cached = self._cache_get(question)
        if cached:
            logger.info("Cache hit")
            _stream_text(cached["answer"], placeholder)
            return {**cached, "cached": True}

        try:
            t0 = time.time()
            source_docs = self._retrieve(question)
            context = "\n\n".join(d.page_content for d in source_docs)

            buf = ""
            with self._claude.messages.stream(
                model=self.cfg.claude_model,
                max_tokens=self.cfg.max_tokens,
                temperature=self.cfg.temperature,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": USER_PROMPT_TEMPLATE.format(
                    context=context, question=question
                )}],
            ) as stream:
                for text in stream.text_stream:
                    buf += text
                    placeholder.markdown(buf + "▌")
                placeholder.markdown(buf)

                final = stream.get_final_message()
                self.metrics.add_usage(
                    final.usage.input_tokens,
                    final.usage.output_tokens,
                )
                logger.info(
                    f"Claude stream — in:{final.usage.input_tokens} "
                    f"out:{final.usage.output_tokens} "
                    f"session_cost:${self.metrics.total_cost:.5f}"
                )

            result = {
                "answer": buf,
                "sources": self._build_sources(source_docs),
                "elapsed": round(time.time() - t0, 1),
                "cached": False,
                "error": False,
            }
            self._cache_set(question, result)
            return result

        except anthropic.APIConnectionError as e:
            msg = "⚠️ Could not reach the Anthropic API. Check your internet connection."
            placeholder.error(msg)
            return {"answer": msg, "sources": [], "elapsed": 0, "cached": False, "error": True}
        except anthropic.RateLimitError:
            msg = "⚠️ Rate limit reached. Please wait a moment and try again."
            placeholder.error(msg)
            return {"answer": msg, "sources": [], "elapsed": 0, "cached": False, "error": True}
        except anthropic.AuthenticationError:
            msg = "⚠️ Invalid API key. Check ANTHROPIC_API_KEY in your .env file."
            placeholder.error(msg)
            return {"answer": msg, "sources": [], "elapsed": 0, "cached": False, "error": True}
        except Exception as e:
            logger.error(f"Streaming query failed: {e}")
            placeholder.error(f"Error: {e}")
            return {"answer": str(e), "sources": [], "elapsed": 0, "cached": False, "error": True}

    # ── Stats ─────────────────────────────────────────────────────────

    def get_stats(self) -> Dict[str, Any]:
        count = 0
        if self.vectorstore:
            try:
                with self._db() as conn:
                    cur = conn.cursor()
                    cur.execute(
                        "SELECT COUNT(*) FROM langchain_pg_embedding e "
                        "JOIN langchain_pg_collection c ON e.collection_id = c.uuid "
                        "WHERE c.name = %s",
                        (self.cfg.collection_name,),
                    )
                    count = cur.fetchone()[0]
            except Exception:
                count = -1
        files = collect_files(self.cfg.docs_path)
        m = self.metrics
        return {
            "chunks_indexed": count,
            "documents_found": len(files),
            "embedding_model": self.cfg.embedding_model,
            "llm_provider": "claude",
            "llm_model": self.cfg.claude_model,
            "chunk_size": self.cfg.chunk_size,
            "retrieval_k": self.cfg.retrieval_k,
            "cache_size": len(self._query_cache),
            "last_indexed": (
                self.last_indexed.strftime("%Y-%m-%d %H:%M") if self.last_indexed else "Never"
            ),
            "session_queries": m.total_queries,
            "session_tokens": m.total_tokens,
            "session_cost_usd": round(m.total_cost, 5),
        }


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def _cb(cb, progress: float, message: str):
    if cb:
        cb(progress, message)
    logger.info(f"[{progress:.0%}] {message}")


def _stream_text(text: str, placeholder, delay: float = 0.012):
    """Fake word-by-word stream for cached responses."""
    words = text.split()
    buf = ""
    for word in words:
        buf += word + " "
        placeholder.markdown(buf + "▌")
        time.sleep(delay)
    placeholder.markdown(buf.strip())


@st.cache_resource(show_spinner=False)
def get_rag_system() -> RAGSystem:
    return RAGSystem()


@st.cache_data(show_spinner=False)
def get_client_config() -> dict:
    return load_client_config()


def _handle_question(rag: RAGSystem, question: str, cfg: dict):
    """Process a question and append to session-state messages."""
    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)
    with st.chat_message("assistant"):
        placeholder = st.empty()
        result = rag.ask_streaming(question, placeholder)   # real streaming
        if not result["error"]:
            ui_cfg = cfg.get("ui", {})
            if ui_cfg.get("show_sources", True) and result["sources"]:
                _render_sources(result["sources"])
            meta = []
            if result.get("cached"):
                meta.append("⚡ cached")
            if ui_cfg.get("show_response_time", True) and result.get("elapsed"):
                meta.append(f"⏱ {result['elapsed']}s")
            if meta:
                st.caption(" · ".join(meta))
            if ui_cfg.get("show_copy_button", True) and result["answer"]:
                st.button(
                    "📋 Copy answer",
                    key=f"copy_{hashlib.md5(question.encode()).hexdigest()[:8]}",
                    on_click=lambda: st.write(
                        f'<script>navigator.clipboard.writeText({json.dumps(result["answer"])})</script>',
                        unsafe_allow_html=True,
                    ),
                )
    st.session_state.messages.append({
        "role": "assistant",
        "content": result["answer"],
        "sources": result.get("sources", []),
        "cached": result.get("cached", False),
        "elapsed": result.get("elapsed"),
    })


def _file_icon(name: str) -> str:
    ext = Path(name).suffix.lower()
    return {"pdf": "📄", "docx": "📝", "doc": "📝", "pptx": "📊", "ppt": "📊",
            "xlsx": "📋", "xls": "📋", "csv": "📋", "txt": "📃"}.get(ext.lstrip("."), "📄")


def _render_sources(sources: List[Dict]):
    if not sources:
        return
    cards_html = ""
    for i, src in enumerate(sources, 1):
        label   = src.get("display_name") or Path(src["metadata"].get("source", "unknown")).name
        fname   = Path(src["metadata"].get("source", label)).name
        icon    = _file_icon(fname)
        preview = src["content"][:180].replace("<", "&lt;").replace(">", "&gt;").replace("\n", " ")
        cards_html += f"""
        <div class="src-card">
          <div class="src-header">
            <span class="src-icon">{icon}</span>
            <span class="src-label">{label}</span>
            <span class="src-num">#{i}</span>
          </div>
          <div class="src-preview">{preview}…</div>
        </div>"""
    with st.expander(f"📚 {len(sources)} source{'s' if len(sources)>1 else ''} referenced"):
        st.markdown(f'<div class="src-grid">{cards_html}</div>', unsafe_allow_html=True)


# ─────────────────────────────────────────────
# Streamlit UI
# ─────────────────────────────────────────────

def run_ui():
    cfg            = get_client_config()
    branding       = cfg.get("branding", {})
    ui_cfg         = cfg.get("ui", {})
    suggested      = cfg.get("suggested_questions", [])
    accent         = branding.get("accent_colour", "#185FA5")
    company        = branding.get("company_name", "RAG Assistant")
    assistant_name = branding.get("assistant_name", "Document Assistant")
    tagline        = branding.get("tagline", "Ask questions across your documents")

    st.set_page_config(
        page_title=f"{company} — {assistant_name}",
        page_icon="🔍",
        layout="wide",
        initial_sidebar_state="expanded",
    )

    # ── Design system ─────────────────────────────────────────────────
    # Derives from accent colour but adds a full modern token set
    st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

    :root {{
        --accent:        {accent};
        --accent-light:  {accent}18;
        --accent-mid:    {accent}40;
        --ink:           #0F172A;
        --ink-mid:       #475569;
        --ink-light:     #94A3B8;
        --surface:       #FFFFFF;
        --surface-2:     #F8FAFC;
        --surface-3:     #F1F5F9;
        --border:        #E2E8F0;
        --border-light:  #F1F5F9;
        --success:       #059669;
        --success-bg:    #ECFDF5;
        --warn:          #D97706;
        --warn-bg:       #FFFBEB;
        --radius-sm:     6px;
        --radius:        10px;
        --radius-lg:     16px;
        --shadow-sm:     0 1px 3px rgba(0,0,0,0.06), 0 1px 2px rgba(0,0,0,0.04);
        --shadow:        0 4px 12px rgba(0,0,0,0.08), 0 2px 4px rgba(0,0,0,0.04);
        --shadow-lg:     0 12px 32px rgba(0,0,0,0.10), 0 4px 8px rgba(0,0,0,0.06);
    }}

    html, body, [class*="css"] {{
        font-family: 'Inter', -apple-system, sans-serif !important;
    }}

    /* ── Sidebar ── */
    [data-testid="stSidebar"] {{
        background: var(--surface-2) !important;
        border-right: 1px solid var(--border) !important;
        min-width: 260px !important;
        max-width: 280px !important;
    }}
    [data-testid="stSidebar"] > div:first-child {{
        padding-top: 0 !important;
        padding-bottom: 1rem !important;
        overflow-y: auto !important;
    }}
    [data-testid="stSidebar"] [data-testid="stVerticalBlock"] {{
        gap: 0.25rem !important;
    }}

    /* ── Brand header ── */
    .brand-header {{
        background: linear-gradient(135deg, var(--accent) 0%, {accent}cc 100%);
        padding: 14px 16px 12px;
        margin: -1rem -1rem 1rem;
        position: relative;
        overflow: hidden;
    }}
    .brand-header::after {{
        content: '';
        position: absolute; top: 0; right: -20px;
        width: 80px; height: 100%;
        background: rgba(255,255,255,0.06);
        transform: skewX(-15deg);
    }}
    .brand-header h1 {{
        color: #fff; font-size: 1rem; font-weight: 700;
        margin: 0 0 2px; letter-spacing: -0.01em;
    }}
    .brand-header p {{
        color: rgba(255,255,255,0.75); font-size: 0.75rem; margin: 0;
    }}
    .brand-dot {{
        display: inline-block; width: 6px; height: 6px;
        border-radius: 50%; background: #4ADE80;
        margin-right: 5px; vertical-align: middle;
        box-shadow: 0 0 6px #4ADE80;
    }}

    /* ── Cost pill ── */
    .cost-pill {{
        background: var(--accent-light);
        border: 1px solid var(--accent-mid);
        border-radius: 999px;
        padding: 3px 10px;
        font-size: 0.76rem; font-weight: 600;
        color: var(--accent);
        display: inline-block; margin-top: 6px;
    }}

    /* ── Sidebar section headers ── */
    .sidebar-section {{
        font-size: 0.65rem; font-weight: 700; letter-spacing: .1em;
        text-transform: uppercase; color: var(--ink-light);
        margin: 12px 0 4px; padding: 0;
    }}

    /* ── File list ── */
    .file-item {{
        display: flex; align-items: center; gap: 8px;
        padding: 6px 8px; border-radius: var(--radius-sm);
        border: 1px solid var(--border-light);
        background: var(--surface);
        margin-bottom: 4px;
        font-size: 0.82rem; color: var(--ink);
    }}
    .file-item .file-name {{ flex: 1; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }}
    .file-item .file-ext {{
        font-size: 0.68rem; font-weight: 700; padding: 1px 5px;
        border-radius: 3px; text-transform: uppercase;
    }}
    .ext-pdf  {{ background: #FEE2E2; color: #991B1B; }}
    .ext-docx, .ext-doc {{ background: #DBEAFE; color: #1E40AF; }}
    .ext-xlsx, .ext-xls, .ext-csv {{ background: #D1FAE5; color: #065F46; }}
    .ext-pptx, .ext-ppt {{ background: #FEF3C7; color: #92400E; }}
    .ext-txt  {{ background: var(--surface-3); color: var(--ink-mid); }}

    /* ── Index buttons ── */
    .index-btn-group {{ display: flex; gap: 6px; margin: 8px 0; }}

    /* ── Stats cards ── */
    .stats-grid {{
        display: grid; grid-template-columns: 1fr 1fr;
        gap: 5px; margin: 6px 0;
    }}
    .stat-card {{
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: var(--radius-sm);
        padding: 6px 8px; text-align: center;
    }}
    .stat-val {{ font-size: 1.1rem; font-weight: 700; color: var(--accent); line-height: 1; }}
    .stat-lbl {{ font-size: 0.65rem; color: var(--ink-light); margin-top: 2px; }}

    /* ── Suggested question pills ── */
    .sq-pill {{
        display: inline-block;
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: 999px;
        padding: 6px 12px;
        font-size: 0.82rem; color: var(--ink);
        cursor: pointer;
        margin: 3px 2px;
        transition: all .15s;
    }}
    .sq-pill:hover {{ background: var(--accent-light); border-color: var(--accent); color: var(--accent); }}

    /* ── Welcome screen ── */
    .welcome-hero {{
        text-align: center;
        padding: 1.5rem 1rem 1rem;
        max-width: 580px;
        margin: 0 auto;
    }}
    .welcome-icon {{
        width: 48px; height: 48px; border-radius: 14px;
        background: linear-gradient(135deg, var(--accent), {accent}99);
        display: flex; align-items: center; justify-content: center;
        font-size: 1.4rem;
        margin: 0 auto 0.75rem;
        box-shadow: var(--shadow-sm);
    }}
    .welcome-hero h2 {{
        font-size: 1.3rem; font-weight: 700; color: var(--ink);
        letter-spacing: -0.02em; margin: 0 0 0.35rem;
    }}
    .welcome-hero p {{
        color: var(--ink-mid); font-size: 0.88rem; line-height: 1.5; margin: 0 0 1rem;
    }}
    .example-grid {{
        display: grid; grid-template-columns: 1fr 1fr;
        gap: 6px; margin-top: 0.75rem; text-align: left;
    }}
    .example-card {{
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: var(--radius);
        padding: 9px 12px;
        font-size: 0.82rem; color: var(--ink-mid);
        cursor: pointer; transition: all .15s;
        line-height: 1.4;
    }}
    .example-card:hover {{ border-color: var(--accent); color: var(--accent); background: var(--accent-light); }}
    .example-card strong {{ display: block; color: var(--ink); margin-bottom: 2px; font-size: 0.8rem; }}

    /* ── Source cards ── */
    .src-grid {{ display: flex; flex-direction: column; gap: 8px; padding: 4px 0; }}
    .src-card {{
        background: var(--surface-2);
        border: 1px solid var(--border);
        border-radius: var(--radius);
        padding: 10px 12px;
        transition: box-shadow .15s;
    }}
    .src-card:hover {{ box-shadow: var(--shadow-sm); }}
    .src-header {{ display: flex; align-items: center; gap: 8px; margin-bottom: 6px; }}
    .src-icon {{ font-size: 1rem; }}
    .src-label {{ font-size: 0.82rem; font-weight: 600; color: var(--ink); flex: 1; }}
    .src-num {{
        font-size: 0.68rem; font-weight: 700;
        background: var(--accent-light); color: var(--accent);
        border-radius: 999px; padding: 1px 7px;
    }}
    .src-preview {{ font-size: 0.78rem; color: var(--ink-mid); line-height: 1.5; }}

    /* ── Step guide cards ── */
    .step-card {{
        border-radius: var(--radius);
        padding: 1rem 1.25rem;
        margin-bottom: 0.75rem;
        border: 1px solid var(--border);
        background: var(--surface);
        box-shadow: var(--shadow-sm);
    }}
    .step-number {{
        display: inline-flex; align-items: center; justify-content: center;
        width: 24px; height: 24px; border-radius: 50%;
        background: var(--accent); color: #fff;
        font-size: 0.72rem; font-weight: 700;
        margin-bottom: 0.5rem;
    }}
    .step-card h3 {{ font-size: 0.95rem; font-weight: 600; color: var(--ink); margin: 0 0 0.2rem; }}
    .step-card p  {{ font-size: 0.84rem; color: var(--ink-mid); margin: 0; }}

    /* ── Chat loading stages ── */
    .thinking-stage {{
        display: flex; align-items: center; gap: 8px;
        font-size: 0.84rem; color: var(--ink-mid);
        padding: 6px 0;
    }}
    .thinking-dot {{
        width: 6px; height: 6px; border-radius: 50%;
        background: var(--accent); opacity: 0.5;
        animation: pulse 1.4s ease-in-out infinite;
    }}
    @keyframes pulse {{ 0%,100% {{ opacity:.3; transform:scale(.8); }} 50% {{ opacity:1; transform:scale(1.2); }} }}

    /* ── General cleanup ── */
    /* Hide Streamlit toolbar and reduce top padding */
    [data-testid="stToolbar"] {{ display: none !important; }}
    header[data-testid="stHeader"] {{ display: none !important; }}
    .block-container {{
        padding-top: 1.5rem !important;
        padding-bottom: 1rem !important;
        max-width: 100% !important;
    }}
    /* Tighten tab content area */
    [data-testid="stTabContent"] {{
        padding-top: 0.5rem !important;
    }}
    /* Tab indicator — match accent colour */
    [data-testid="stTabs"] [role="tablist"] {{
        border-bottom: 1px solid var(--border) !important;
        gap: 0 !important;
    }}
    [data-testid="stTabs"] button[role="tab"] {{
        font-size: 0.88rem !important;
        font-weight: 500 !important;
        color: var(--ink-mid) !important;
        border-bottom: 2px solid transparent !important;
        padding: 0.5rem 1rem !important;
    }}
    [data-testid="stTabs"] button[role="tab"][aria-selected="true"] {{
        color: var(--accent) !important;
        border-bottom-color: var(--accent) !important;
        font-weight: 600 !important;
    }}
    /* Hide the default red BaseWeb tab highlight */
    [data-baseweb="tab-highlight"] {{ display: none !important; }}
    .stButton > button {{
        border-radius: var(--radius-sm) !important;
        font-weight: 500 !important;
        font-size: 0.85rem !important;
        transition: all .15s !important;
    }}
    .stTextInput > div > div > input {{
        border-radius: var(--radius-sm) !important;
        border-color: var(--border) !important;
    }}
    .stTextInput > div > div > input:focus {{
        border-color: var(--accent) !important;
        box-shadow: 0 0 0 3px {accent}20 !important;
    }}
    [data-testid="stChatMessage"] {{
        padding: 0.5rem 0 !important;
    }}
    .stTabs [data-baseweb="tab"] {{
        font-size: 0.88rem !important;
    }}
    div[data-testid="stExpander"] {{
        border-radius: var(--radius) !important;
        border-color: var(--border) !important;
    }}
    </style>
    """, unsafe_allow_html=True)

    rag = get_rag_system()

    # ── Auto-connect if everything is ready ──────────────────────────
    # No "Connect" button needed — if key + index exist, connect silently
    files = collect_files(config.docs_path)  # needed for auto-connect check
    if not rag._initialized:
        fresh_key = os.getenv("ANTHROPIC_API_KEY", "") or load_saved_api_key()
        if fresh_key and files:
            _auto_chunks = 0
            try:
                _ac = psycopg2.connect(
                    host=config.db_host, port=int(config.db_port),
                    dbname=config.db_name, user=config.db_user,
                    password=config.db_password,
                )
                _acur = _ac.cursor()
                _acur.execute(
                    "SELECT COUNT(*) FROM langchain_pg_embedding e "
                    "JOIN langchain_pg_collection c ON e.collection_id = c.uuid "
                    "WHERE c.name = %s", (config.collection_name,)
                )
                _auto_chunks = _acur.fetchone()[0]
                _ac.close()
            except Exception:
                pass
            if _auto_chunks > 0:
                # Silently connect in background — user sees chat immediately
                with st.spinner("Opening your knowledge base..."):
                    rag.setup(rebuild=False)

    # ── Sidebar ──────────────────────────────────────────────────────
    with st.sidebar:
        st.markdown(f"""
        <div class="brand-header">
            <h1>{company}</h1>
            <p><span class="brand-dot"></span>{assistant_name}</p>
        </div>
        """, unsafe_allow_html=True)

        # ── API key / Configuration ───────────────────────────────────
        with st.expander("⚙️ Configuration", expanded=not (os.getenv("ANTHROPIC_API_KEY","") or load_saved_api_key())):
            st.caption(f"**Model:** `{config.claude_model}`")
            saved_key = load_saved_api_key()
            env_key   = os.getenv("ANTHROPIC_API_KEY", "")
            has_key   = bool(env_key or saved_key)

            if env_key:
                st.success("✅ API key configured")
            elif saved_key:
                st.success("✅ API key saved")
                if st.button("Remove key", key="del_key"):
                    delete_api_key()
                    config.anthropic_api_key = ""
                    rag._initialized = False
                    rag._claude = None
                    st.rerun()
            else:
                st.warning("No API key set")
                new_key = st.text_input(
                    "Anthropic API key",
                    type="password",
                    placeholder="sk-ant-...",
                    key="api_key_input",
                )
                if st.button("Save key", key="save_key", type="primary"):
                    if new_key.startswith("sk-ant-"):
                        save_api_key(new_key)
                        config.anthropic_api_key = new_key
                        st.rerun()
                    else:
                        st.error("Key must start with sk-ant-")
                st.caption("[Get your key →](https://console.anthropic.com)")

            if has_key and rag._initialized:
                m = rag.metrics
                st.markdown(
                    f'<span class="cost-pill">'
                    f'💰 ${m.total_cost:.4f} · {m.total_tokens:,} tokens · {m.total_queries} queries'
                    f'</span>',
                    unsafe_allow_html=True,
                )

        st.markdown("---")

        # ── Knowledge base ────────────────────────────────────────────
        files = collect_files(config.docs_path)
        chunks_in_db = 0
        try:
            tmp_conn = psycopg2.connect(
                host=config.db_host, port=int(config.db_port),
                dbname=config.db_name, user=config.db_user,
                password=config.db_password,
            )
            cur = tmp_conn.cursor()
            cur.execute(
                "SELECT COUNT(*) FROM langchain_pg_embedding e "
                "JOIN langchain_pg_collection c ON e.collection_id = c.uuid "
                "WHERE c.name = %s",
                (config.collection_name,),
            )
            chunks_in_db = cur.fetchone()[0]
            tmp_conn.close()
            logger.info(f"Sidebar DB check: {chunks_in_db} chunks in collection '{config.collection_name}'")
        except Exception as e:
            logger.warning(f"Sidebar DB check failed: {e}")

        kb_status = "🟢 Ready" if (chunks_in_db > 0 and rag._initialized) else ("🟡 Not connected" if chunks_in_db > 0 else "🔴 Not indexed")
        st.markdown(f'<p class="sidebar-section">Knowledge Base &nbsp; <span style="font-size:0.75rem;text-transform:none;font-weight:500;letter-spacing:0">{kb_status}</span></p>', unsafe_allow_html=True)

        if files:
            file_html = ""
            for f in files:
                ext = f.suffix.lower().lstrip(".")
                icon = _file_icon(f.name)
                file_html += f'''
                <div class="file-item">
                  <span>{icon}</span>
                  <span class="file-name">{f.name}</span>
                  <span class="file-ext ext-{ext}">{ext}</span>
                </div>'''
            with st.expander(f"📂 {len(files)} document{'s' if len(files)>1 else ''} · {chunks_in_db:,} chunks", expanded=False):
                st.markdown(file_html, unsafe_allow_html=True)
        else:
            st.warning("No documents found")

        # ── Change detection ──────────────────────────────────────────
        if chunks_in_db > 0:
            _rag_for_hash = rag if rag.vectorstore else RAGSystem()
            stored_hashes = _rag_for_hash._load_hash_store()
            if stored_hashes:
                stored_by_name = {Path(k).name: v for k, v in stored_hashes.items()}
                new_or_changed = [f.name for f in files if stored_by_name.get(f.name) != _rag_for_hash._get_file_hash(f)]
                deleted = set(Path(k).name for k in stored_hashes) - set(f.name for f in files)
                if new_or_changed or deleted:
                    changed_list = ", ".join(list(new_or_changed)[:2])
                    if len(new_or_changed) > 2:
                        changed_list += f" +{len(new_or_changed)-2} more"
                    st.warning(f"⚠️ {len(new_or_changed)} new/changed file(s) — click Refresh")
        elif files:
            st.warning("⚠️ Documents not indexed — click Build Index")

        st.markdown("---")

        # ── Action buttons ────────────────────────────────────────────
        st.markdown('<p class="sidebar-section">Actions</p>', unsafe_allow_html=True)

        if not rag._initialized:
            init_btn = st.button("▶ Open Assistant", use_container_width=True, type="primary",
                                  help="Connect to your document index and start chatting")
        else:
            init_btn = False

        col1, col2 = st.columns(2)
        with col1:
            update_btn  = st.button("↻ Refresh", use_container_width=True,
                                     help="Re-index new or changed documents only")
        with col2:
            rebuild_btn = st.button("⊕ Build Index", use_container_width=True,
                                     help="Rebuild the entire document index from scratch")

        st.markdown("---")

        # ── Stats ─────────────────────────────────────────────────────
        if rag._initialized:
            stats = rag.get_stats()
            with st.expander("📊 Activity", expanded=False):
                st.markdown(f'''
                <div class="stats-grid">
                  <div class="stat-card">
                    <div class="stat-val">{stats["chunks_indexed"]}</div>
                    <div class="stat-lbl">Chunks</div>
                  </div>
                  <div class="stat-card">
                    <div class="stat-val">{stats["documents_found"]}</div>
                    <div class="stat-lbl">Docs</div>
                  </div>
                  <div class="stat-card">
                    <div class="stat-val">{stats["session_queries"]}</div>
                    <div class="stat-lbl">Questions</div>
                  </div>
                  <div class="stat-card">
                    <div class="stat-val">{stats["cache_size"]}</div>
                    <div class="stat-lbl">Cached</div>
                  </div>
                </div>
                ''', unsafe_allow_html=True)
                st.caption(f"Last indexed: {stats['last_indexed']}")
                if st.button("Clear cache", use_container_width=True):
                    rag._query_cache.clear()
                    st.success("Cache cleared")

        # Suggested questions shown in main chat area — not duplicated in sidebar

        # ── Trust signals ─────────────────────────────────────────────
        st.markdown('''
        <div style="padding:10px 0 4px">
          <div style="font-size:0.75rem;color:#64748B;line-height:1.8">
            <div>✓ &nbsp;Documents stay on your server</div>
            <div>✓ &nbsp;Powered by Claude AI</div>
            <div>✓ &nbsp;Your data is never used for training</div>
          </div>
          <div style="font-size:0.7rem;color:#94A3B8;margin-top:8px">
            Powered by <strong>WeldAI</strong> · weldai.uk
          </div>
        </div>
        ''', unsafe_allow_html=True)

    # ── Init / Rebuild handler ────────────────────────────────────────
    if init_btn or rebuild_btn or update_btn:
        rebuild_mode = "incremental" if update_btn else (True if rebuild_btn else False)
        progress_bar = st.progress(0.0)
        status_txt   = st.empty()

        def on_progress(frac: float, msg: str):
            progress_bar.progress(min(frac, 1.0))
            status_txt.text(msg)

        with st.spinner("Building your knowledge base..."):
            ok = rag.setup(rebuild=rebuild_mode, progress_cb=on_progress)

        progress_bar.empty()
        status_txt.empty()
        if ok:
            st.success("✅ Knowledge base ready — start asking questions!")
            st.rerun()
        else:
            st.error("❌ Setup failed. Check that your API key is valid and documents are readable.")

    # ── Not yet initialised — guide user ─────────────────────────────
    if not rag._initialized:
        fresh_key = os.getenv("ANTHROPIC_API_KEY", "") or load_saved_api_key()

        if not fresh_key:
            st.markdown(f'''
            <div class="welcome-hero">
              <div class="welcome-icon">🔑</div>
              <h2>Welcome to {company}</h2>
              <p>{tagline}</p>
            </div>''', unsafe_allow_html=True)
            st.markdown('''<div class="step-card">
              <div class="step-number">1</div>
              <h3>Connect your AI</h3>
              <p>Open <strong>⚙️ Configuration</strong> in the sidebar and enter your Anthropic API key.
              Get your key at <a href="https://console.anthropic.com" target="_blank">console.anthropic.com</a></p>
            </div>''', unsafe_allow_html=True)
            return

        if not files:
            st.markdown(f'''
            <div class="welcome-hero">
              <div class="welcome-icon">📂</div>
              <h2>Add your documents</h2>
              <p>Copy your business documents into the docs folder on the server, then refresh this page.</p>
            </div>''', unsafe_allow_html=True)
            st.markdown('''<div class="step-card">
              <div class="step-number">2</div>
              <h3>Load your knowledge base</h3>
              <p>Supported: <strong>PDF · Word · Excel · PowerPoint · CSV · Text</strong></p>
            </div>''', unsafe_allow_html=True)
            return

        # Has key + docs — check if indexed
        _chunks2 = chunks_in_db  # already computed above in sidebar
        example_qs = suggested[:4] if suggested else [
            "What does this document cover?",
            "Summarise the key points",
            "What are the main steps?",
            "List the key findings",
        ]
        cards = "".join([f'<div class="example-card">💬 {q}</div>' for q in example_qs])

        if _chunks2 > 0:
            st.markdown(f'''
            <div class="welcome-hero" style="padding-top:1.5rem">
              <div class="welcome-icon">🔍</div>
              <h2>Ready to search</h2>
              <p style="margin-bottom:0.5rem">{_chunks2:,} knowledge chunks · {len(files)} documents · {tagline}</p>
              <div class="example-grid">{cards}</div>
            </div>''', unsafe_allow_html=True)
            st.markdown('''<div class="step-card" style="max-width:480px;margin:1rem auto">
              <div class="step-number">▶</div>
              <h3>Click "Open Assistant" to start chatting</h3>
              <p>Press the <strong>▶ Open Assistant</strong> button in the sidebar.</p>
            </div>''', unsafe_allow_html=True)
        else:
            st.markdown(f'''
            <div class="welcome-hero" style="padding-top:1.5rem">
              <div class="welcome-icon">⚡</div>
              <h2>{len(files)} document(s) ready to index</h2>
              <p>Your files are loaded. Build the index once and you can ask questions in seconds.</p>
              <div class="example-grid">{cards}</div>
            </div>''', unsafe_allow_html=True)
            st.markdown('''<div class="step-card" style="max-width:480px;margin:1rem auto">
              <div class="step-number">▶</div>
              <h3>Click "Build Index" to get started</h3>
              <p>Press <strong>⊕ Build Index</strong> in the sidebar. Takes about 1 minute for a typical document set.</p>
            </div>''', unsafe_allow_html=True)
        return

    # ── Main tabs ─────────────────────────────────────────────────────
    tab_chat, tab_search, tab_about = st.tabs(["💬 Chat", "🔎 Search", "ℹ️ About"])

    # ── Chat ─────────────────────────────────────────────────────────
    with tab_chat:
        if "messages" not in st.session_state:
            st.session_state.messages = []

        pending = st.session_state.pop("pending_question", None)

        # ── Empty state ────────────────────────────────────────────────
        if not st.session_state.messages and not pending:
            stats2 = rag.get_stats()
            example_qs2 = suggested[:4] if suggested else [
                "What does this document cover?",
                "Summarise the key points",
                "What are the main steps?",
                "List the key findings",
            ]
            cards2 = "".join([
                f'<div class="example-card" onclick="void(0)">💬 {q}</div>'
                for q in example_qs2
            ])
            st.markdown(f'''
            <div style="max-width:580px;margin:0.25rem auto 0.25rem;text-align:center;padding:0 0.5rem">
              <p style="font-size:0.7rem;font-weight:700;letter-spacing:.1em;text-transform:uppercase;
                color:var(--ink-light);margin:0 0 0.3rem">
                {stats2["chunks_indexed"]:,} chunks &nbsp;·&nbsp; {stats2["documents_found"]} documents &nbsp;·&nbsp;
                <span style="color:#059669">🟢 Ready</span>
              </p>
              <h2 style="font-size:1.1rem;font-weight:700;color:var(--ink);
                letter-spacing:-0.02em;margin:0 0 0.2rem">
                Ask anything about your documents
              </h2>
              <p style="font-size:0.8rem;color:var(--ink-mid);margin:0 0 0.5rem;line-height:1.4">
                Instant answers with exact source and page number cited.
              </p>
              <div class="example-grid">{cards2}</div>
            </div>
            ''', unsafe_allow_html=True)

        if prompt := st.chat_input(f"Ask {assistant_name} anything about your documents..."):
            _handle_question(rag, prompt, cfg)
            st.rerun()

        for msg in reversed(st.session_state.messages):
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])
                if msg["role"] == "assistant":
                    if ui_cfg.get("show_sources", True) and msg.get("sources"):
                        _render_sources(msg["sources"])
                    meta = []
                    if msg.get("cached"):
                        meta.append("⚡ cached")
                    if ui_cfg.get("show_response_time", True) and msg.get("elapsed"):
                        meta.append(f"⏱ {msg['elapsed']}s")
                    if meta:
                        st.caption(" · ".join(meta))

        if pending:
            _handle_question(rag, pending, cfg)
            st.rerun()

        if st.session_state.messages:
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🗑️ Clear conversation"):
                    st.session_state.messages = []
                    st.rerun()
            with col2:
                last_user = next(
                    (m["content"] for m in reversed(st.session_state.messages) if m["role"] == "user"),
                    None,
                )
                if last_user and st.button("🔁 Ask again"):
                    _handle_question(rag, last_user, cfg)
                    st.rerun()

    # ── Search ────────────────────────────────────────────────────────
    with tab_search:
        st.markdown("##### 🔎 Direct document search")
        st.caption("Search your document index directly — useful for checking what's been indexed or finding specific passages.")
        query = st.text_input("Search query", placeholder="e.g. refund policy, safety procedures, invoice terms...")
        k = st.slider("Number of results", 1, 10, 4)
        if st.button("Search documents", type="primary") and query:
            with st.spinner("Searching..."):
                docs = rag.vectorstore.similarity_search(query, k=k)
            st.caption(f"Found {len(docs)} matching passages")
            for i, doc in enumerate(docs, 1):
                src   = Path(doc.metadata.get("source", "unknown")).name
                page  = doc.metadata.get("page", "")
                label = f"{src} — p.{int(page)+1}" if page != "" else src
                icon  = _file_icon(src)
                with st.expander(f"{icon} {label}"):
                    st.markdown(f'''<div style="font-size:0.88rem;line-height:1.6;color:#334155;
                        background:#F8FAFC;padding:12px;border-radius:8px;border:1px solid #E2E8F0">
                        {doc.page_content.replace(chr(10), "<br>")}
                        </div>''', unsafe_allow_html=True)

    # ── About ─────────────────────────────────────────────────────────
    with tab_about:
        stats3 = rag.get_stats()
        st.markdown(f"##### About {company}")
        col1, col2 = st.columns(2)
        with col1:
            st.markdown(f"""
**Product:** {assistant_name}
**AI Model:** Claude API (`{config.claude_model}`)
**Embeddings:** Local (fastembed — no API cost)
**Vector store:** PostgreSQL + pgvector
            """)
        with col2:
            st.markdown(f"""
**Documents:** {stats3["documents_found"]} files
**Index size:** {stats3["chunks_indexed"]:,} chunks
**Last indexed:** {stats3["last_indexed"]}
**Session queries:** {stats3["session_queries"]}
            """)
        st.markdown("---")
        st.markdown("**Supported document formats**")
        st.markdown("PDF · Word (.docx .doc) · PowerPoint (.pptx .ppt) · Excel (.xlsx .xls) · CSV · Plain text · RTF")
        with st.expander("Technical details"):
            st.json(stats3)


# ─────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────

def cli_index():
    print(f"[{datetime.now()}] Indexing '{config.docs_path}'...")
    rag = RAGSystem()
    rag.init_embeddings()
    rag.connect_vectorstore()
    n = rag.index_documents()
    print(f"[{datetime.now()}] Done — {n} chunks indexed.")

def cli_query(question: str):
    rag = RAGSystem()
    rag.setup()
    result = rag.ask(question)
    print(f"\nAnswer:\n{result['answer']}\n")
    for i, src in enumerate(result["sources"], 1):
        print(f"Source {i}: {src.get('display_name', '?')}")

def cli_status():
    rag = RAGSystem()
    rag.init_embeddings()
    rag.connect_vectorstore()
    print(json.dumps(rag.get_stats(), indent=2))

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="RAG Pipeline — Claude Edition")
    parser.add_argument("--mode", choices=["ui", "index", "query", "status"], default="ui")
    parser.add_argument("--query", type=str, default="")
    args = parser.parse_args()
    if args.mode == "index":
        cli_index()
    elif args.mode == "query":
        if not args.query:
            print("--query is required"); sys.exit(1)
        cli_query(args.query)
    elif args.mode == "status":
        cli_status()
    else:
        run_ui()